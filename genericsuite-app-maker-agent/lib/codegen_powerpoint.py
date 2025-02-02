"""
PowerPoint generation
"""
import os

# Reference:
# pip install python-pptx
# https://python-pptx.readthedocs.io/en/latest/user/quickstart.html

import pptx
# from pptx.util import Inches

from lib.codegen_utilities import (
    create_dirs,
    read_file,
    log_debug,
)

DEBUG = False

# DEFAULT_POWERPOINT_TEMPLATE = "default_powerpoint_template.pptx"
DEFAULT_POWERPOINT_TEMPLATE = ""


class PowerPointGenerator:
    """
    PowerPoint generator class
    """
    def __init__(self, params: dict = None):
        self.params = params or {}

    def generate(self, slides_config: list):
        """
        Generates the PowerPoint slides
        """
        output_dir = self.params.get("output_dir", "./output")
        file_name = self.params.get("file_name", "app_presentation")
        create_dirs(output_dir)
        target_file_path = f"{output_dir}/{file_name}.pptx"

        template = self.params.get("template", DEFAULT_POWERPOINT_TEMPLATE)
        if template:
            template_path = f"./config/{template}"
            if not os.path.exists(template_path):
                raise ValueError(f"Missing template file: {template_path}")
            pptx_template = read_file(template_path)
            pptx_obj = pptx.Presentation(pptx_template)
        else:
            pptx_obj = pptx.Presentation()

        bullet_slide_layout = pptx_obj.slide_layouts[1]
        for slide_config in slides_config.get("slides", []):
            slide = pptx_obj.slides.add_slide(bullet_slide_layout)
            shapes = slide.shapes

            log_debug("PowerPointGenerator | generate | "
                      f"slide_config: {slide_config}", debug=DEBUG)

            # shapes.title.font.size = Inches(2)
            title_shape = shapes.title
            body_shape = shapes.placeholders[1]

            title_shape.text = slide_config.get("title")
            content_items = slide_config.get("content")
            if isinstance(content_items, str):
                content_items = [{
                    "type": "text",
                    "text": content_items,
                }]
            for content in content_items:
                if content.get("type") == "text":
                    # text = shapes.add_textbox(
                    #     # Inches(content.get("x")),
                    #     # Inches(content.get("y")),
                    #     # Inches(content.get("width")),
                    #     # Inches(content.get("height")),
                    # )
                    # First bullet text
                    tf = body_shape.text_frame
                    text = content.get("text", "")
                    if "\n" in text or "\r" in text or "* " in text:
                        separator = \
                            "\n" if "\n" in text else \
                            "* " if "* " in text else "\r"
                        splitted_text = text.split(separator)
                        for line in splitted_text:
                            p = tf.add_paragraph()
                            p.text = line
                            p.level = 0
                    else:
                        tf.text = text
                    # # Use _TextFrame.text for first bullet'
                    # p = tf.add_paragraph()
                    # p.text = "Sub-bullets"
                    # p.level = 1
                    # # Use _TextFrame.add_paragraph() for subsequent bullets
                    # p = tf.add_paragraph()
                    # p.text = 'Subsequent bullets'
                    # p.level = 2

                elif content.get("type") == "image":
                    image = shapes.add_picture(
                        # Inches(content.get("x")),
                        # Inches(content.get("y")),
                        # Inches(content.get("width")),
                        # Inches(content.get("height")),
                        content.get("image_path"),
                    )
                    # image.width = Inches(content.get("width"))
                    # image.height = Inches(content.get("height"))
                    image.shapes.title.text = content.get("title")
                elif content.get("type") == "table":
                    table = shapes.add_table(
                        # Inches(content.get("x")),
                        # Inches(content.get("y")),
                        # Inches(content.get("width")),
                        # Inches(content.get("height")),
                    )
                    for row in content.get("rows", []):
                        for cell in row:
                            table.cell(row=row.get("row"),
                                       column=cell.get("column")).text = \
                                cell.get("text")
                # else:
                #     raise ValueError(
                #       f"Invalid content type: {content.get('type')}")

            # Add notes
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = ''
            if slide_config.get("speaker_notes"):
                text_frame.text += slide_config.get("speaker_notes")
            if slide_config.get("image_prompt"):
                text_frame.text += "\n\nImage Prompt: " + \
                                   slide_config.get("image_prompt")

        # Save presentation and return the file path
        pptx_obj.save(target_file_path)
        return target_file_path
